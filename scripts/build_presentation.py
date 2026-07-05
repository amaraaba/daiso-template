"""Build DAISO presentation (PPTX) with large fonts and a native-shape architecture diagram.

Run:  python scripts/build_presentation.py
Output: DAISO_Presentation.pptx in the repo root.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- Theme ----------
INTEL_BLUE = RGBColor(0x00, 0x68, 0xB5)
DARK = RGBColor(0x1F, 0x2A, 0x44)
GREY = RGBColor(0x55, 0x5E, 0x6E)
LIGHT_BG = RGBColor(0xF3, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xFF, 0x8A, 0x00)
GREEN = RGBColor(0x2E, 0xA0, 0x4B)

FONT = "Calibri"
FOOTER_TEXT = "DAISO  •  Ahmad Maraaba  •  May 2026"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_text(slide, x, y, w, h, text, *, size=24, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=22, color=DARK, bullet="•  "):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = f"{bullet}{item}"
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_footer(slide, page_num, total=8):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SH - Inches(0.45),
                                 SW, Inches(0.45))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = INTEL_BLUE
    add_text(slide, Inches(0.4), SH - Inches(0.42), Inches(8), Inches(0.4),
             FOOTER_TEXT, size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SW - Inches(2), SH - Inches(0.42), Inches(1.6), Inches(0.4),
             f"{page_num} / {total}", size=12, color=WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_header_band(slide, title):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.1))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT_BG
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                    Inches(0.18), Inches(1.1))
    accent.line.fill.background()
    accent.fill.solid()
    accent.fill.fore_color.rgb = INTEL_BLUE
    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
             title, size=32, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)


def new_slide(title=None, page=None):
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    if title:
        add_header_band(slide, title)
    if page is not None:
        add_footer(slide, page)
    return slide


def add_arrow(slide, x1, y1, x2, y2, color=GREY, width=2.5):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")


# ============================================================
# Slide 1 — Title
# ============================================================
s = prs.slides.add_slide(BLANK)
s.background.fill.solid()
s.background.fill.fore_color.rgb = DARK

block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.5), SH)
block.line.fill.background()
block.fill.solid()
block.fill.fore_color.rgb = INTEL_BLUE

add_text(s, Inches(1), Inches(1.4), Inches(11.5), Inches(1.0),
         "DAISO", size=72, bold=True, color=ACCENT)
add_text(s, Inches(1), Inches(2.5), Inches(11.5), Inches(0.6),
         "DTSE AI Solution Orchestrator", size=28, color=WHITE)
add_text(s, Inches(1), Inches(3.7), Inches(11.5), Inches(1.6),
         "From One Copilot to a Shared DTSE Tool",
         size=44, bold=True, color=WHITE)
add_text(s, Inches(1), Inches(5.1), Inches(11.5), Inches(0.8),
         "“You saw mine. Here's how we make it everyone's.”",
         size=22, italic=True, color=RGBColor(0xC8, 0xD2, 0xE0))
add_text(s, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
         "Ahmad Maraaba  •  May 2026", size=18, color=RGBColor(0xC8, 0xD2, 0xE0))


# ============================================================
# Slide 2 — Why a Shared Tool
# ============================================================
s = new_slide("Why a Shared Tool", page=2)

add_text(s, Inches(0.6), Inches(1.4), Inches(12.2), Inches(1.4),
         "“AI assistants are only as good as the knowledge behind them.\n"
         "Knowledge isolated per engineer is knowledge wasted.”",
         size=24, italic=True, color=GREY)

add_bullets(s, Inches(0.8), Inches(3.3), Inches(12), Inches(3),
            ["Within a team — engineers can't share the skills and knowledge they each capture",
             "Across teams — common capabilities (spec search, vault lookup, debug patterns) get rebuilt instead of reused",
             "Result today — every engineer starts from zero, every team reinvents the same wiring"],
            size=22)

band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.85))
band.line.fill.background()
band.fill.solid()
band.fill.fore_color.rgb = INTEL_BLUE
add_text(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.85),
         "DAISO makes knowledge contribute-once, reuse-everywhere",
         size=24, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# Slide 3 — DAISO Structure
# ============================================================
s = new_slide("DAISO — One Repo, Federated by Team", page=3)

add_bullets(s, Inches(0.8), Inches(1.5), Inches(12), Inches(2.8),
            ["dtse-common/  —  skills and knowledge everyone shares",
             "packs/<team>/  —  team-owned domain  (Array, Scan, Functional, Reset)",
             "packs/<team>/products/<product>/  —  silicon-program overlays  (NVL, PTL, ARL)"],
            size=22)

def outcome_box(x, title, body, color):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             x, Inches(4.6), Inches(5.8), Inches(2.0))
    box.line.color.rgb = color
    box.line.width = Pt(2)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    add_text(s, x + Inches(0.3), Inches(4.7), Inches(5.4), Inches(0.5),
             title, size=22, bold=True, color=color)
    add_text(s, x + Inches(0.3), Inches(5.3), Inches(5.4), Inches(1.4),
             body, size=18, color=DARK)

outcome_box(Inches(0.8), "Engineers get",
            "Their team's pack +\nevery shared skill, automatically.", INTEL_BLUE)
outcome_box(Inches(6.8), "Engineers give back",
            "Skills and knowledge they capture\nflow into the pack for the whole team.", GREEN)


# ============================================================
# Slide 4 — Architecture (native shapes)
# ============================================================
s = new_slide("Architecture", page=4)

def arch_box(x, y, w, h, label, fill, text_color=WHITE, size=20):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    add_text(s, x, y, w, h, label, size=size, bold=True,
             color=text_color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return shp

arch_box(Inches(0.6), Inches(2.0), Inches(3.7), Inches(1.7),
         "Windows VS Code\n(git source of truth)", INTEL_BLUE)
add_text(s, Inches(0.6), Inches(3.8), Inches(3.7), Inches(0.4),
         "engineer edits here", size=14, italic=True,
         color=GREY, align=PP_ALIGN.CENTER)

add_text(s, Inches(4.4), Inches(2.3), Inches(1.6), Inches(0.5),
         "robocopy", size=16, italic=True, color=GREY, align=PP_ALIGN.CENTER)
add_text(s, Inches(4.4), Inches(2.85), Inches(1.6), Inches(0.5),
         "(sync.ps1)", size=13, italic=True, color=GREY, align=PP_ALIGN.CENTER)

arch_box(Inches(6.1), Inches(2.0), Inches(3.5), Inches(1.7),
         "NFS\n(MCP runtime)", DARK)
add_text(s, Inches(6.1), Inches(3.8), Inches(3.5), Inches(0.4),
         "MCP server runs here", size=14, italic=True,
         color=GREY, align=PP_ALIGN.CENTER)

arch_box(Inches(10.0), Inches(2.0), Inches(2.8), Inches(1.7),
         "GitHub\nCopilot", GREEN)
add_text(s, Inches(10.0), Inches(3.8), Inches(2.8), Inches(0.4),
         "in VS Code chat", size=14, italic=True,
         color=GREY, align=PP_ALIGN.CENTER)

# arrows
add_arrow(s, Inches(4.3), Inches(2.85), Inches(6.1), Inches(2.85))
add_arrow(s, Inches(9.6), Inches(2.6), Inches(10.0), Inches(2.6))
add_arrow(s, Inches(10.0), Inches(3.1), Inches(9.6), Inches(3.1))
add_text(s, Inches(9.4), Inches(3.3), Inches(1.2), Inches(0.4),
         "SSH stdio", size=13, italic=True, color=GREY, align=PP_ALIGN.CENTER)

note = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.6), Inches(5.3), Inches(12.2), Inches(1.4))
note.line.color.rgb = INTEL_BLUE
note.line.width = Pt(1.5)
note.fill.solid()
note.fill.fore_color.rgb = LIGHT_BG
add_text(s, Inches(0.9), Inches(5.45), Inches(11.6), Inches(0.5),
         "Git stays invisible to engineers", size=22, bold=True, color=DARK)
add_text(s, Inches(0.9), Inches(5.95), Inches(11.6), Inches(0.6),
         "scripts/sync.ps1 handles pull + sync to NFS in one command.",
         size=18, color=GREY)


# ============================================================
# Slide 5 — Onboarding
# ============================================================
s = new_slide('Onboarding — "Set Me Up"', page=5)

col_w = Inches(5.8)
row_h = Inches(1.5)
left_x = Inches(0.6)
right_x = Inches(6.9)
top_y = Inches(1.5)

hdr_left = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, top_y, col_w, Inches(0.7))
hdr_left.fill.solid(); hdr_left.fill.fore_color.rgb = GREY; hdr_left.line.fill.background()
add_text(s, left_x, top_y, col_w, Inches(0.7), "Before",
         size=24, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

hdr_right = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, top_y, col_w, Inches(0.7))
hdr_right.fill.solid(); hdr_right.fill.fore_color.rgb = INTEL_BLUE; hdr_right.line.fill.background()
add_text(s, right_x, top_y, col_w, Inches(0.7), "Now",
         size=24, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

rows = [
    ("Per-engineer init prompt\n+ manual MCP wiring",
     'Clone repo → say "set me up"\n→ guided 15 min'),
    ("Each engineer maintains\ntheir own files",
     "Personal config only\n(mcp.json, sync conf)"),
]
y = top_y + Inches(0.8)
for i, (lt, rt) in enumerate(rows):
    bg = WHITE if i % 2 == 0 else LIGHT_BG
    for x, txt in [(left_x, lt), (right_x, rt)]:
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_w, row_h)
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        cell.line.color.rgb = RGBColor(0xDD, 0xE3, 0xEB)
        cell.line.width = Pt(0.75)
        add_text(s, x + Inches(0.25), y, col_w - Inches(0.5), row_h, txt,
                 size=22, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    y += row_h


# ============================================================
# Slide 6 — What's in the Box Today
# ============================================================
s = new_slide("What's in the Box Today", page=6)

add_bullets(s, Inches(0.7), Inches(1.5), Inches(12.2), Inches(5.2),
            ["4 packs scaffolded:  Array (full), Scan, Functional, Reset",
             "Array pack:  7 skills,  32 NVL knowledge files,  5 MCP tools (ported from DFT Copilot)",
             "dtse-common:  10 shared skills  (setup, contribute, discover, search-specs, vault-helper, glossary, compare, session-summary, status-report, explain-file)",
             'Invisible git workflow:  engineer says "save my changes" → branch + commit + push happens silently'],
            size=22)


# ============================================================
# Slide 7 — How Knowledge Flows Back
# ============================================================
s = new_slide("How Knowledge Flows Back", page=7)

steps = [
    ("1", "Engineer hits\na new failure"),
    ("2", "AI captures\nthe fix"),
    ("3", 'Engineer says\n"save my changes"'),
    ("4", "Goes into team's\npack"),
    ("5", "Next engineer\ngets it on sync"),
]
n = len(steps)
box_w = Inches(2.3)
box_h = Inches(1.8)
gap = Inches(0.15)
total_w = box_w * n + gap * (n - 1)
start_x = (SW - total_w) // 2
y = Inches(2.0)

for i, (num, label) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    color = INTEL_BLUE if i < n - 1 else GREEN
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    add_text(s, x, y + Inches(0.1), box_w, Inches(0.5), num,
             size=30, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.1), y + Inches(0.75), box_w - Inches(0.2), Inches(1.0),
             label, size=15, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < n - 1:
        ax1 = x + box_w
        ax2 = ax1 + gap
        ay = y + box_h // 2
        add_arrow(s, ax1, ay, ax2, ay, color=DARK, width=2)

promo = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.0))
promo.fill.solid(); promo.fill.fore_color.rgb = ACCENT
promo.line.fill.background()
add_text(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.0),
         "If broadly useful  →  promote to dtse-common  →  every team benefits",
         size=22, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.8),
         "The loop that compounds: more use → more knowledge → sharper tool",
         size=20, italic=True, color=GREY, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 8 — The Ask
# ============================================================
s = new_slide("The Ask", page=8)

add_text(s, Inches(0.8), Inches(1.6), Inches(12), Inches(1.0),
         "Each team:  nominate a pack champion",
         size=32, bold=True, color=DARK)

add_bullets(s, Inches(1.2), Inches(2.9), Inches(11.5), Inches(2.6),
            ["Owns their team's pack — reviews contributions, sets domain rules",
             "Onboards engineers as they come on",
             "First point of contact for their team's pack"],
            size=22)

box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.1))
box.fill.solid(); box.fill.fore_color.rgb = INTEL_BLUE
box.line.fill.background()
add_text(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.1),
         "Who's the champion for your pack?",
         size=28, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


out = "DAISO_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
