"""Build a single 'Contributing Changes' slide for the DAISO presentation.

Run:  python scripts/build_contribute_slide.py
Output: DAISO_Contribute_Slide.pptx in the repo root.

Theme matches scripts/build_presentation.py so you can copy the slide
straight into your existing deck.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- Theme (same as build_presentation.py) ----------
INTEL_BLUE = RGBColor(0x00, 0x68, 0xB5)
DARK = RGBColor(0x1F, 0x2A, 0x44)
GREY = RGBColor(0x55, 0x5E, 0x6E)
LIGHT_BG = RGBColor(0xF3, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xFF, 0x8A, 0x00)
GREEN = RGBColor(0x2E, 0xA0, 0x4B)
CODE_BG = RGBColor(0x1F, 0x2A, 0x44)
CODE_FG = RGBColor(0xE6, 0xED, 0xF3)
CODE_COMMENT = RGBColor(0x8B, 0x96, 0xA8)

FONT = "Calibri"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_text(slide, x, y, w, h, text, *, size=24, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
             font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_code_block(slide, x, y, w, h, lines, *, size=14):
    """lines: list of (text, is_comment) tuples."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.fill.background()

    tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15),
                                  w - Inches(0.4), h - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (line, is_comment) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = MONO
        r.font.size = Pt(size)
        r.font.color.rgb = CODE_COMMENT if is_comment else CODE_FG


def add_header_band(slide, title):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.1))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT_BG
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                    Inches(0.18), Inches(1.1))
    accent.line.fill.background()
    accent.fill.solid()
    accent.fill.fore_color.rgb = INTEL_BLUE
    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
             title, size=32, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SH - Inches(0.45),
                                 SW, Inches(0.45))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = INTEL_BLUE
    add_text(slide, Inches(0.4), SH - Inches(0.42), Inches(8), Inches(0.4),
             "DAISO  •  Ahmad Maraaba  •  June 2026",
             size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# Slide — Contributing Changes
# ============================================================
s = prs.slides.add_slide(BLANK)
s.background.fill.solid()
s.background.fill.fore_color.rgb = WHITE
add_header_band(s, "Contributing Changes")
add_footer(s)

# ---- LEFT COLUMN: Manual Way (git) ----
left_x = Inches(0.6)
col_w = Inches(6.2)

add_text(s, left_x, Inches(1.25), col_w, Inches(0.5),
         "The Manual Way  —  git",
         size=22, bold=True, color=GREY)

code_lines = [
    ("# 1. Create a branch", True),
    ("git checkout -b you/short-desc", False),
    ("", False),
    ("# 2. Stage only files you want to share", True),
    ("git add packs/array/skills/my-skill/...", False),
    ("", False),
    ("# 3. Commit with a message", True),
    ('git commit -m "Add my-skill"', False),
    ("", False),
    ("# 4. Push to remote", True),
    ("git push -u origin you/short-desc", False),
    ("", False),
    ("# 5. Switch back to main", True),
    ("git checkout main", False),
]
add_code_block(s, left_x, Inches(1.8), col_w, Inches(3.7), code_lines, size=14)

# Rules under code
rules_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               left_x, Inches(5.65), col_w, Inches(1.25))
rules_box.line.color.rgb = ACCENT
rules_box.line.width = Pt(1.5)
rules_box.fill.solid()
rules_box.fill.fore_color.rgb = LIGHT_BG
add_text(s, left_x + Inches(0.25), Inches(5.85), col_w - Inches(0.5), Inches(0.45),
         "Rules", size=16, bold=True, color=ACCENT)
add_text(s, left_x + Inches(0.25), Inches(6.2), col_w - Inches(0.5), Inches(0.7),
         "•  Never push to main — always branch\n"
         "•  Branch:  idsid/short-desc   •   Never  git add -A",
         size=13, color=DARK)

# ---- RIGHT COLUMN: Easy Way ----
right_x = Inches(7.0)
right_w = Inches(5.8)

add_text(s, right_x, Inches(1.25), right_w, Inches(0.5),
         "The Easy Way  —  Copilot",
         size=22, bold=True, color=INTEL_BLUE)

# Big trigger pill
pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          right_x, Inches(1.85), right_w, Inches(0.85))
pill.fill.solid()
pill.fill.fore_color.rgb = INTEL_BLUE
pill.line.fill.background()
add_text(s, right_x, Inches(1.85), right_w, Inches(0.85),
         '"push my changes"',
         size=26, bold=True, color=WHITE, italic=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Steps
add_text(s, right_x, Inches(2.95), right_w, Inches(0.4),
         "What Copilot does for you:", size=15, bold=True, color=DARK)

steps = [
    "1.  Lists what changed — plain English, no git jargon",
    "2.  Asks which files to share  (by number / all / skip private)",
    "3.  Branches, commits, pushes — only what you picked",
    "4.  Hands you the PR link to click",
]
tb = s.shapes.add_textbox(right_x, Inches(3.4), right_w, Inches(2.0))
tf = tb.text_frame
tf.word_wrap = True
for i, line in enumerate(steps):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6)
    r = p.add_run()
    r.text = line
    r.font.name = FONT
    r.font.size = Pt(15)
    r.font.color.rgb = DARK

# Other triggers
trig_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              right_x, Inches(5.4), right_w, Inches(1.5))
trig_box.line.color.rgb = GREEN
trig_box.line.width = Pt(1.5)
trig_box.fill.solid()
trig_box.fill.fore_color.rgb = LIGHT_BG
add_text(s, right_x + Inches(0.25), Inches(5.45), right_w - Inches(0.5), Inches(0.4),
         "Other phrases that work", size=15, bold=True, color=GREEN)
add_text(s, right_x + Inches(0.25), Inches(5.85), right_w - Inches(0.5), Inches(1.0),
         '"save my changes"   •   "submit"\n'
         '"I\'m done"   •   "share my work"',
         size=15, italic=True, color=DARK)

# Bottom callout — PR link
pr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.6), Inches(6.95) - Inches(0.5) - Inches(0.05),
                        SW - Inches(1.2), Inches(0.0))
# (omitted bottom band — we keep it clean above the footer)

out = "DAISO_Contribute_Slide.pptx"
prs.save(out)
print(f"Wrote {out}")
