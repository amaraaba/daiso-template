"""TFM lifecycle / where-it-runs slide (Idea 1: swimlane pipeline).

A single slide with:
  - Time columns left -> right (AUTHOR, DEPOSIT, GEN-LIST, RUN, INGEST, RELEASE)
  - Lanes top -> bottom: Engineer / TFM / Runner (simless, sim, emu) / Storage / MPE
  - Save points drop into the Storage lane
  - Runner lane fans into 3 sub-lanes only inside the RUN column

Run:    python scripts/build_tfm_flow_slide.py
Output: TFM_Flow_Slide.pptx in repo root.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---- Theme ----
INTEL_BLUE = RGBColor(0x00, 0x68, 0xB5)
DARK = RGBColor(0x1F, 0x2A, 0x44)
GREY = RGBColor(0x55, 0x5E, 0x6E)
LIGHT = RGBColor(0xF3, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xFF, 0x8A, 0x00)
GREEN = RGBColor(0x2E, 0xA0, 0x4B)
LANE_BG_A = RGBColor(0xF7, 0xFA, 0xFE)
LANE_BG_B = RGBColor(0xEE, 0xF3, 0xFA)
LANE_LINE = RGBColor(0xC7, 0xD4, 0xE5)
COL_HDR_BG = RGBColor(0x1F, 0x2A, 0x44)
COL_HDR_FG = RGBColor(0xFF, 0xFF, 0xFF)
DB_FILL = RGBColor(0xE6, 0xED, 0xF3)

FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_text(slide, x, y, w, h, text, *, size=12, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def card(slide, x, y, w, h, title, fill, *, sub=None, text_color=WHITE,
         size=13, sub_size=10, line_color=None,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1.0)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = text_color
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sub
        r2.font.name = FONT
        r2.font.size = Pt(sub_size)
        r2.font.bold = False
        r2.font.color.rgb = text_color
    return shp


def line(slide, x1, y1, x2, y2, *, color=DARK, width=2.0, dashed=False,
         arrow=True, kind=MSO_CONNECTOR_TYPE.STRAIGHT):
    conn = slide.shapes.add_connector(kind, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    if dashed:
        d = etree.SubElement(ln, qn("a:prstDash"))
        d.set("val", "dash")
    if arrow:
        tail = etree.SubElement(ln, qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("h", "med")
    conn.shadow.inherit = False
    return conn


# ============================================================
# Slide
# ============================================================
s = prs.slides.add_slide(BLANK)
s.background.fill.solid()
s.background.fill.fore_color.rgb = WHITE

# Header
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.85))
band.line.fill.background()
band.fill.solid()
band.fill.fore_color.rgb = LIGHT
band.shadow.inherit = False
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), Inches(0.85))
acc.line.fill.background()
acc.fill.solid()
acc.fill.fore_color.rgb = INTEL_BLUE
acc.shadow.inherit = False
add_text(s, Inches(0.40), Inches(0.08), Inches(9.5), Inches(0.42),
         "TFM Lifecycle  —  where the test runs and where the data lives",
         size=24, bold=True, color=DARK)
add_text(s, Inches(0.42), Inches(0.50), Inches(9.5), Inches(0.30),
         "From pattern_spec to MPE release  ·  simless / sim / emu fan-out  ·  save points highlighted",
         size=12, italic=True, color=GREY)

# ---------- grid geometry ----------
LANE_LBL_X = Inches(0.20)
LANE_LBL_W = Inches(1.55)
GRID_X = LANE_LBL_X + LANE_LBL_W + Inches(0.05)
GRID_W = SW - GRID_X - Inches(0.30)
N_COLS = 6
COL_W = Emu(int(GRID_W / N_COLS))

COL_HDR_Y = Inches(0.95)
COL_HDR_H = Inches(0.42)

LANES_Y = COL_HDR_Y + COL_HDR_H + Inches(0.05)

# Lane definitions: (key, label, height_in)
LANES = [
    ("eng",     "Engineer",          0.65),
    ("tfm",     "TFM",               0.85),
    ("simless", "Runner — simless",  0.60),
    ("sim",     "Runner — sim",      0.60),
    ("emu",     "Runner — emu",      0.60),
    ("store",   "Storage",           0.85),
    ("mpe",     "MPE / HVM",         0.55),
]

# compute lane Y positions
y_cursor = LANES_Y
lane_y = {}
lane_h = {}
GAP = Inches(0.04)
for i, (k, _, h) in enumerate(LANES):
    lane_y[k] = y_cursor
    lane_h[k] = Inches(h)
    y_cursor = y_cursor + Inches(h) + GAP

# Column headers (time)
COL_TITLES = ["AUTHOR", "DEPOSIT TEST", "GEN-LIST", "RUN",
              "INGEST STIL", "RELEASE"]
for i, t in enumerate(COL_TITLES):
    cx = GRID_X + COL_W * i
    cw = COL_W - Inches(0.06)
    hdr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(0.03),
                             COL_HDR_Y, cw, COL_HDR_H)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = COL_HDR_BG
    hdr.line.fill.background()
    hdr.shadow.inherit = False
    tf = hdr.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = t
    r.font.name = FONT
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = COL_HDR_FG

# Lane label column + lane backgrounds
for i, (k, lbl, _) in enumerate(LANES):
    y = lane_y[k]
    h = lane_h[k]
    # lane label
    add_text(s, LANE_LBL_X, y, LANE_LBL_W, h, lbl,
             size=12, bold=True, color=DARK,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    # lane background (alternating)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, GRID_X, y, GRID_W, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LANE_BG_A if i % 2 == 0 else LANE_BG_B
    bg.line.color.rgb = LANE_LINE
    bg.line.width = Pt(0.5)
    bg.shadow.inherit = False
    # send to back so cards overlay it
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)

# helper: get a cell rectangle
def cell(lane_key, col_idx, *, pad=0.10):
    x = GRID_X + COL_W * col_idx + Inches(pad)
    w = COL_W - Inches(pad * 2)
    y = lane_y[lane_key] + Inches(0.06)
    h = lane_h[lane_key] - Inches(0.12)
    return x, y, w, h


def cell_center(lane_key, col_idx):
    x, y, w, h = cell(lane_key, col_idx)
    return x + w / 2, y + h / 2


# ---------- Cards ----------
# Engineer
x, y, w, h = cell("eng", 0)
card(s, x, y, w, h, "pattern_spec", DARK, sub="Tessent + Tim", size=12)

# TFM (cols 1, 2, 4, 5)
for ci, (title, sub) in [
    (1, ("deposit test", "TID + version")),
    (2, ("gen-list", "turnin → .list")),
    (4, ("deposit stil", "--from-run")),
    (5, ("tag golden", "MPE handoff")),
]:
    x, y, w, h = cell("tfm", ci)
    card(s, x, y, w, h, title, AMBER, sub=sub, size=13, sub_size=10)

# Runner sub-lanes (col 3 = RUN)
for k, title, sub in [
    ("simless", "pattern → .stil", "fast · no sim"),
    ("sim",     "VCS sim → .stil + waves", "full RTL sim"),
    ("emu",     "ZeBu emu → .stil + run_history", "fastest on big patterns"),
]:
    x, y, w, h = cell(k, 3)
    card(s, x, y, w, h, title, GREEN, sub=sub, size=12, sub_size=9.5)

# Storage row — DB cylinders at save points
def db(slide, x, y, w, h, title, sub):
    return card(slide, x, y, w, h, title, DB_FILL, sub=sub,
                text_color=DARK, size=12, sub_size=10,
                line_color=GREY, shape=MSO_SHAPE.CAN)

# Tests DB at DEPOSIT col
x, y, w, h = cell("store", 1)
db(s, x, y, w, h, "Tests DB", "test_record")

# NFS rundir at RUN col
x, y, w, h = cell("store", 3)
db(s, x, y, w, h, "NFS rundir", ".stil + logs")

# STILs DB + Apparate at INGEST col — split the cell horizontally
x, y, w, h = cell("store", 4)
half_w = (w - Inches(0.10)) / 2
db(s, x, y, half_w, h, "STILs DB", "stil_record\nunique_id")
db(s, x + half_w + Inches(0.10), y, half_w, h, "Apparate", "object store")

# MPE
x, y, w, h = cell("mpe", 5)
card(s, x, y, w, h, "MPE / HVM", DARK, sub="release consumer", size=12, sub_size=10)

# ---------- Connectors ----------
# Engineer.AUTHOR -> TFM.deposit
ax, ay = cell_center("eng", 0)
tx, ty = cell_center("tfm", 1)
line(s, ax + Inches(0.55), ay,
     tx - Inches(0.55), ty - Inches(0.05),
     color=DARK, width=2.0)

# TFM main horizontal pipeline (deposit -> gen-list -> [ingest] -> release)
tfm_y = lane_y["tfm"] + lane_h["tfm"] / 2
# deposit -> gen-list
x1, _ = cell_center("tfm", 1)
x2, _ = cell_center("tfm", 2)
line(s, x1 + Inches(0.55), tfm_y, x2 - Inches(0.55), tfm_y,
     color=AMBER, width=2.5)
# gen-list -> deposit stil (skip RUN col, but show as faded handoff arrow)
x2, _ = cell_center("tfm", 2)
x4, _ = cell_center("tfm", 4)
line(s, x2 + Inches(0.55), tfm_y, x4 - Inches(0.55), tfm_y,
     color=AMBER, width=2.5, dashed=True)
add_text(s, (x2 + x4) / 2 - Inches(1.0), tfm_y - Inches(0.30),
         Inches(2.0), Inches(0.24),
         "(handoff via Runner)", size=10, italic=True, color=GREY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# deposit stil -> tag golden
x4, _ = cell_center("tfm", 4)
x5, _ = cell_center("tfm", 5)
line(s, x4 + Inches(0.55), tfm_y, x5 - Inches(0.55), tfm_y,
     color=AMBER, width=2.5)

# gen-list -> Runner fan-out (3 arrows down to each run mode)
gx, _ = cell_center("tfm", 2)
gen_bottom_y = lane_y["tfm"] + lane_h["tfm"]
# vertical drop to a junction below TFM
junction_x = GRID_X + COL_W * 3 + COL_W / 2  # middle of RUN col
junction_y = gen_bottom_y + Inches(0.05)
line(s, gx, gen_bottom_y, junction_x, junction_y,
     color=GREEN, width=1.75, arrow=False)
# from junction down/over to each run mode card
for k in ("simless", "sim", "emu"):
    rx, ry = cell_center(k, 3)
    rcell_x, _, rcell_w, _ = cell(k, 3)
    line(s, junction_x, junction_y,
         rcell_x, ry, color=GREEN, width=1.75, dashed=True)
add_text(s, junction_x - Inches(1.5), junction_y + Inches(0.02),
         Inches(3.0), Inches(0.22),
         ".list  →  Runner",
         size=10, bold=True, color=GREEN,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

# Runner outputs -> deposit stil (3 fan-in arrows going up-right)
ix, iy = cell_center("tfm", 4)
deposit_left_x = ix - Inches(0.55)
for k in ("simless", "sim", "emu"):
    rcx, rcy = cell_center(k, 3)
    rcell_x, _, rcell_w, _ = cell(k, 3)
    line(s, rcell_x + rcell_w, rcy,
         deposit_left_x, iy, color=GREEN, width=1.5, dashed=True)
add_text(s, deposit_left_x - Inches(2.4), iy - Inches(0.55),
         Inches(2.2), Inches(0.22),
         ".stil + run_result  ↑",
         size=10, bold=True, color=GREEN,
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP)

# Save-point arrows: TFM step -> Storage cell directly below
def save_arrow(tfm_col, store_col, color=AMBER):
    sx, _ = cell_center("tfm", tfm_col)
    s_top_y = lane_y["tfm"] + lane_h["tfm"]
    stx, sty = cell_center("store", store_col)
    line(s, sx, s_top_y, stx, lane_y["store"],
         color=color, width=2.0, dashed=False)

save_arrow(1, 1)              # deposit test -> Tests DB
save_arrow(4, 4)              # deposit stil -> STILs DB / Apparate column

# Run-mode -> NFS rundir (one arrow from middle of run lanes to rundir)
mid_run_y = (lane_y["sim"] + lane_h["sim"] / 2)
runx, _ = cell_center("store", 3)
rundir_top_y = lane_y["store"]
line(s, runx, mid_run_y + Inches(0.50), runx, rundir_top_y,
     color=GREEN, width=1.5, dashed=True)

# Storage RELEASE -> MPE
sdx, sdy = cell_center("store", 4)
mx, my = cell_center("mpe", 5)
# Use orange/release arrow from STILs DB area down-right to MPE
line(s, GRID_X + COL_W * 5, sdy,
     mx - Inches(0.55), my, color=DARK, width=2.2, dashed=False)
add_text(s, GRID_X + COL_W * 5 - Inches(0.10), sdy + Inches(0.05),
         Inches(2.0), Inches(0.22),
         "release",
         size=10, bold=True, color=DARK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

# tag golden -> MPE (also)
tgx, tgy = cell_center("tfm", 5)
line(s, tgx, lane_y["tfm"] + lane_h["tfm"], mx, lane_y["mpe"],
     color=AMBER, width=1.75, dashed=True)

# ---------- Bottom takeaway ----------
bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(0.30), SH - Inches(0.65),
                         Inches(12.73), Inches(0.50))
bar.line.fill.background()
bar.fill.solid()
bar.fill.fore_color.rgb = INTEL_BLUE
bar.shadow.inherit = False
add_text(s, Inches(0.30), SH - Inches(0.65), Inches(12.73), Inches(0.50),
         "TFM owns the rails (deposit · gen-list · ingest · release)  ·  "
         "the Runner is plugged in via .list / .stil  ·  "
         "data drops into Tests DB, STILs DB, and Apparate at fixed save points.",
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

# Legend (top-right, small)
LX = Inches(10.10)
for i, (lab, col) in enumerate([("Engineer / external", DARK),
                                ("TFM step", AMBER),
                                ("Runner step", GREEN),
                                ("Storage", GREY)]):
    cy = Inches(0.10) + Inches(0.20) * i
    chip = s.shapes.add_shape(MSO_SHAPE.OVAL, LX, cy, Inches(0.17), Inches(0.17))
    chip.fill.solid()
    chip.fill.fore_color.rgb = col
    chip.line.fill.background()
    chip.shadow.inherit = False
    add_text(s, LX + Inches(0.24), cy - Inches(0.02), Inches(2.9), Inches(0.22),
             lab, size=10, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

out = "TFM_Flow_Slide.pptx"
prs.save(out)
print(f"Saved {out}")
